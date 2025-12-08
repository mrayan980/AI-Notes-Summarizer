import re
from pathlib import Path

def extract_text_from_file(filepath):
    """
    Extract text from various file formats
    """
    file_extension = Path(filepath).suffix.lower()
    
    if file_extension == '.pdf':
        return extract_from_pdf(filepath)
    elif file_extension in ['.ppt', '.pptx']:
        return extract_from_ppt(filepath)
    elif file_extension == '.txt':
        return extract_from_txt(filepath)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def extract_from_pdf(filepath):
    """
    Extract text from PDF using PyPDF2
    """
    try:
        import PyPDF2
        
        text = []
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                
                if page_text.strip():
                    text.append(f"\n--- Page {page_num + 1} ---\n")
                    text.append(page_text)
        
        extracted = ''.join(text)
        return clean_text(extracted)
    
    except ImportError:
        # Fallback to pdfplumber if PyPDF2 is not available
        try:
            import pdfplumber
            
            text = []
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    
                    if page_text and page_text.strip():
                        text.append(f"\n--- Page {page_num + 1} ---\n")
                        text.append(page_text)
            
            extracted = ''.join(text)
            return clean_text(extracted)
        
        except ImportError:
            raise ImportError("Please install PyPDF2 or pdfplumber: pip install PyPDF2 pdfplumber")

def extract_from_ppt(filepath):
    """
    Extract text from PowerPoint files using python-pptx
    """
    try:
        from pptx import Presentation
        
        text = []
        prs = Presentation(filepath)
        
        for slide_num, slide in enumerate(prs.slides):
            text.append(f"\n--- Slide {slide_num + 1} ---\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
                    text.append('\n')
        
        extracted = ''.join(text)
        return clean_text(extracted)
    
    except ImportError:
        raise ImportError("Please install python-pptx: pip install python-pptx")

def extract_from_txt(filepath):
    """
    Extract text from plain text files
    """
    try:
        with open(filepath, 'r', encoding='utf-8') as file:
            text = file.read()
        return clean_text(text)
    
    except UnicodeDecodeError:
        # Try with different encoding
        with open(filepath, 'r', encoding='latin-1') as file:
            text = file.read()
        return clean_text(text)

def clean_text(text):
    """
    Clean and normalize extracted text
    """
    # Remove excessive whitespace
    text = re.sub(r'\n\s*\n', '\n\n', text)
    
    # Remove excessive spaces
    text = re.sub(r' +', ' ', text)
    
    # Remove special characters that cause issues
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    
    # Strip leading/trailing whitespace
    text = text.strip()
    
    return text

def split_into_sentences(text):
    """
    Split text into sentences
    """
    # Simple sentence splitting
    sentences = re.split(r'(?<=[.!?])\s+', text)
    return [s.strip() for s in sentences if s.strip()]

def split_into_lines(text):
    """
    Split text into lines
    """
    lines = text.split('\n')
    return [line.strip() for line in lines if line.strip()]

def get_text_metadata(text):
    """
    Extract metadata about the text
    """
    return {
        'word_count': len(text.split()),
        'char_count': len(text),
        'line_count': len(text.split('\n')),
        'sentence_count': len(split_into_sentences(text))
    }